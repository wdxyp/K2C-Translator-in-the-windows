import warnings
# 静默 openpyxl 关于 DrawingML 支持不全的警告
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl.reader.drawings")

import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
import pandas as pd
import os
import json
import hashlib
import random
import string
import requests
from openpyxl import load_workbook
from openpyxl.styles import Font, Border, Alignment, PatternFill
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.oxml.ns import nsmap as DOCX_NSMAP
from datetime import datetime
import re
import time
import threading
import shutil

# 请配置 KIMI API KEY（建议通过环境变量注入，避免提交到仓库）
KIMI_API_KEY = os.getenv('KIMI_API_KEY') or os.getenv('MOONSHOT_API_KEY') or 'YOUR_KIMI_API_KEY'
KIMI_MODEL = os.getenv('KIMI_MODEL') or 'kimi-k2.5'
API_URL = os.getenv('KIMI_API_URL') or 'https://api.moonshot.cn/v1/chat/completions'

def _load_local_kimi_credentials():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    for filename in ("kimi_credentials.json", "kimi_credentials.local.json"):
        path = os.path.join(base_dir, filename)
        if not os.path.exists(path):
            continue
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f) or {}
            api_key = str(data.get("KIMI_API_KEY") or data.get("MOONSHOT_API_KEY") or "").strip()
            model = str(data.get("KIMI_MODEL") or "").strip()
            api_url = str(data.get("KIMI_API_URL") or "").strip()
            if api_key:
                return api_key, model, api_url
        except Exception:
            continue
    return None, None, None

if KIMI_API_KEY == "YOUR_KIMI_API_KEY":
    _api_key, _model, _api_url = _load_local_kimi_credentials()
    if _api_key:
        KIMI_API_KEY = _api_key
        if _model:
            KIMI_MODEL = _model
        if _api_url:
            API_URL = _api_url

# ==========================================
# V2.13 (KIMI API)
# ==========================================

# 全局变量
revision_map = {}
original_texts = []
translated_texts = []
HANGUL_RE = re.compile(r"[\uac00-\ud7a3]")
_kimi_request_lock = threading.Lock()
_kimi_rate_lock = threading.Lock()
_kimi_next_allowed_ts = 0.0
_kimi_min_interval = 20.0
_translation_cache = {}
_last_progress_ts = 0.0
_last_progress_msg = ""
_last_api_ts = 0.0
_last_api_status = ""
_watchdog_stop_event = threading.Event()
_watchdog_thread = None
_last_watchdog_ping_ts = 0.0

def _kimi_throttle():
    global _kimi_next_allowed_ts
    with _kimi_rate_lock:
        now = time.monotonic()
        wait = _kimi_next_allowed_ts - now
        if wait > 0:
            time.sleep(wait)
            now = time.monotonic()
        _kimi_next_allowed_ts = now + _kimi_min_interval

def _kimi_post(payload, timeout):
    with _kimi_request_lock:
        _kimi_throttle()
        headers = {
            'Authorization': f'Bearer {KIMI_API_KEY}',
            'Content-Type': 'application/json',
        }
        return requests.post(API_URL, headers=headers, json=payload, timeout=timeout)

def _mark_progress(msg):
    global _last_progress_ts, _last_progress_msg
    _last_progress_ts = time.time()
    _last_progress_msg = str(msg)

def _mark_api(status):
    global _last_api_ts, _last_api_status
    _last_api_ts = time.time()
    _last_api_status = str(status)

def _watchdog_loop():
    global _last_watchdog_ping_ts
    while not _watchdog_stop_event.is_set():
        time.sleep(20)
        now = time.time()
        last_progress = _last_progress_ts
        last_api = _last_api_ts
        msg = _last_progress_msg or "运行中"
        progress_age = int(now - last_progress) if last_progress else -1
        api_age = int(now - last_api) if last_api else -1
        status = _last_api_status or "unknown"
        if progress_age >= 60 or api_age >= 60:
            print(f"[心跳] {msg} | 距上次进度 {progress_age}s | 距上次API {api_age}s | 上次API {status}")
        if api_age >= 180 and (now - _last_watchdog_ping_ts) >= 300 and KIMI_API_KEY and KIMI_API_KEY != "YOUR_KIMI_API_KEY":
            _last_watchdog_ping_ts = now
            try:
                payload = {
                    'model': KIMI_MODEL,
                    'messages': [
                        {'role': 'user', 'content': 'Translate Korean to Chinese: 안녕하세요'},
                    ],
                    'temperature': 1,
                }
                r = _kimi_post(payload, timeout=15)
                _mark_api(f"http={r.status_code}")
                print(f"[心跳] KIMI 联通性测试: http={r.status_code}")
            except Exception as e:
                _mark_api(f"exception={type(e).__name__}")
                print(f"[心跳] KIMI 联通性测试异常: {type(e).__name__}")

def _watchdog_start():
    global _watchdog_thread
    _watchdog_stop_event.clear()
    _watchdog_thread = threading.Thread(target=_watchdog_loop, daemon=True)
    _watchdog_thread.start()

def _watchdog_stop():
    _watchdog_stop_event.set()

FROM_LANG_MAP = {
    'zh2ko': 'zh', 'ko2zh': 'kor', 'ko2en': 'kor', 'zh2en': 'zh', 'en2zh': 'en',
    'zh_tw2en': 'zh', 'en2zh_tw': 'en', 'zh2ja': 'zh', 'ja2zh': 'ja',
    'en2ko': 'en', 'vi2zh': 'vie', 'zh2vi': 'zh', 'ko2vi': 'kor'
}
TO_LANG_MAP = {
    'zh2ko': 'kor', 'ko2zh': 'zh', 'ko2en': 'en', 'zh2en': 'en', 'en2zh': 'zh',
    'zh_tw2en': 'en', 'en2zh_tw': 'zh', 'zh2ja': 'ja', 'ja2zh': 'zh',
    'en2ko': 'kor', 'vi2zh': 'zh', 'zh2vi': 'vie', 'ko2vi': 'vie'
}
TARGET_FONT_BY_TO_LANG = {'zh': '微软雅黑', 'kor': 'Malgun Gothic'}
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def get_target_to_lang():
    try:
        direction = translation_direction.get()
    except Exception:
        return None
    return TO_LANG_MAP.get(direction)

def get_target_font_name():
    to_lang = get_target_to_lang()
    if not to_lang:
        return None
    return TARGET_FONT_BY_TO_LANG.get(to_lang)

def xpath_with_ns(element, expr):
    if element is None:
        return []
    try:
        return element.xpath(expr, namespaces=DOCX_NSMAP)
    except TypeError:
        return element.xpath(expr)

def set_docx_r_element_font(r_element, font_name):
    if r_element is None or not font_name:
        return
    rPr = r_element.get_or_add_rPr()
    rFonts = rPr.rFonts
    if rFonts is None:
        rFonts = OxmlElement('w:rFonts')
        rPr.append(rFonts)
    rFonts.set(qn('w:ascii'), font_name)
    rFonts.set(qn('w:hAnsi'), font_name)
    rFonts.set(qn('w:eastAsia'), font_name)
    rFonts.set(qn('w:cs'), font_name)
    for k in (qn('w:asciiTheme'), qn('w:hAnsiTheme'), qn('w:eastAsiaTheme'), qn('w:csTheme')):
        try:
            rFonts.attrib.pop(k, None)
        except Exception:
            pass

def set_docx_run_font(run, font_name):
    if not run or not font_name:
        return
    run.font.name = font_name
    set_docx_r_element_font(run._r, font_name)

def set_drawingml_rpr_element_font(a_rPr_element, font_name):
    if a_rPr_element is None or not font_name:
        return
    for local_name, tag in (('latin', 'a:latin'), ('ea', 'a:ea'), ('cs', 'a:cs')):
        el = a_rPr_element.find(f'{{{A_NS}}}{local_name}')
        if el is None:
            el = OxmlElement(tag)
            a_rPr_element.append(el)
        el.set('typeface', font_name)

def set_drawingml_r_element_font(a_r_element, font_name):
    if a_r_element is None or not font_name:
        return
    a_rPr = a_r_element.find(f'{{{A_NS}}}rPr')
    if a_rPr is None:
        a_rPr = OxmlElement('a:rPr')
        a_r_element.insert(0, a_rPr)
    set_drawingml_rpr_element_font(a_rPr, font_name)

def set_pptx_run_font(run, font_name):
    if not run or not font_name:
        return
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r = run._r
        if hasattr(r, 'get_or_add_rPr'):
            rPr = r.get_or_add_rPr()
        else:
            rPr = r.find(f'{{{A_NS}}}rPr')
            if rPr is None:
                rPr = OxmlElement('a:rPr')
                r.insert(0, rPr)
        set_drawingml_rpr_element_font(rPr, font_name)
    except Exception:
        return

def set_pptx_paragraph_default_font(paragraph, font_name):
    if not paragraph or not font_name:
        return
    try:
        p = paragraph._p
        pPr = p.find(f'{{{A_NS}}}pPr')
        if pPr is None:
            pPr = OxmlElement('a:pPr')
            p.insert(0, pPr)
        defRPr = pPr.find(f'{{{A_NS}}}defRPr')
        if defRPr is None:
            defRPr = OxmlElement('a:defRPr')
            pPr.append(defRPr)
        set_drawingml_rpr_element_font(defRPr, font_name)
    except Exception:
        return

def load_revision_dict(file_path="revision.md", silent=False):
    """
    自动加载校准词典 (revision.md)
    """
    mapping = {}
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line or line.startswith("#"): continue
                    if not re.match(r'^[-*+]\s*', line): continue
                    content = re.sub(r'^[-*+]\s*', '', line)
                    if "格式" in content: continue
                    
                    if ":" in content:
                        parts = content.split(":", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
                    elif "：" in content:
                        parts = content.split("：", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
            if not silent:
                print(f"[系统] 自动加载校准文件成功，共 {len(mapping)} 条有效规则。")
        except Exception as e:
            print(f"[错误] 加载校准文件失败: {e}")
    return mapping

def apply_revisions(text):
    """
    在翻译结果上应用校准映射
    """
    if not text or not revision_map: return text
    sorted_keys = sorted(revision_map.keys(), key=len, reverse=True)
    result_text = text
    for err in sorted_keys:
        if err in result_text:
            result_text = result_text.replace(err, revision_map[err])
    return result_text

def kimi_translate(q, from_lang, to_lang):
    if not q or not q.strip():
        return q
    if not KIMI_API_KEY or KIMI_API_KEY == "YOUR_KIMI_API_KEY":
        return q
    cache_key = (from_lang, to_lang, str(q))
    cached = _translation_cache.get(cache_key)
    if cached is not None:
        return cached
    max_retries = 8
    try:
        env_max_retries = os.getenv("KIMI_MAX_RETRIES")
        if env_max_retries is not None and str(env_max_retries).strip() != "":
            max_retries = max(1, int(str(env_max_retries).strip()))
    except Exception:
        max_retries = 8
    retry_base_delay = 2.0
    retry_max_delay = 60.0
    retries = 0
    temperature = 1
    try:
        env_temp = os.getenv("KIMI_TEMPERATURE")
        if env_temp is not None and str(env_temp).strip() != "":
            temperature = float(str(env_temp).strip())
    except Exception:
        temperature = 1
    original_stripped = str(q).strip()
    strict_no_hangul = (from_lang == "kor") and (to_lang != "kor") and bool(HANGUL_RE.search(original_stripped))
    user_prompt_template = (
        f"Translate from {from_lang} to {to_lang}.\n"
        "Return ONLY the translated text.\n"
        "Text:\n"
        "{text}"
    )
    while retries < max_retries:
        try:
            _mark_api("start")
            system_prompt = (
                "You are a translation engine.\n"
                "Rules:\n"
                "1) Output ONLY the translation.\n"
                "2) Do NOT add explanations, notes, examples, pinyin, brackets, quotes, markdown, or extra text.\n"
                "3) Preserve the original line breaks EXACTLY (same count and positions).\n"
                "4) Preserve punctuation and spacing style as much as possible.\n"
            )
            user_prompt = user_prompt_template.format(text=q)
            payload = {
                'model': KIMI_MODEL,
                'messages': [
                    {'role': 'system', 'content': system_prompt},
                    {'role': 'user', 'content': user_prompt},
                ],
                'temperature': temperature,
            }
            response = _kimi_post(payload, timeout=60)
            response.encoding = 'utf-8'
            if response.status_code != 200:
                _mark_api(f"http={response.status_code}")
                if response.status_code in (401, 403):
                    return q
                if response.status_code == 400:
                    try:
                        body = response.json() or {}
                        err = body.get("error") or {}
                        msg = str(err.get("message") or "")
                        if ("invalid temperature" in msg) and ("only 1 is allowed" in msg) and payload.get("temperature") != 1:
                            payload["temperature"] = 1
                            response = _kimi_post(payload, timeout=60)
                            response.encoding = 'utf-8'
                            if response.status_code == 200:
                                _mark_api("http=200")
                                result = response.json() or {}
                                choices = result.get('choices') or []
                                if choices:
                                    message = (choices[0] or {}).get('message') or {}
                                    content = message.get('content')
                                    if content is not None:
                                        translated = str(content).strip()
                                        if strict_no_hangul and (translated == original_stripped or HANGUL_RE.search(translated)):
                                            retries += 1
                                            continue
                                        _translation_cache[cache_key] = translated
                                        return translated
                    except Exception:
                        pass
                retries += 1
                if retries < max_retries:
                    delay = min(retry_base_delay * (2 ** (retries - 1)), retry_max_delay)
                    delay = delay + (random.random() * 0.5)
                    time.sleep(delay)
                    continue
                return q
            result = response.json() or {}
            if 'error' in result:
                _mark_api("json_error")
                retries += 1
                if retries < max_retries:
                    delay = min(retry_base_delay * (2 ** (retries - 1)), retry_max_delay)
                    delay = delay + (random.random() * 0.5)
                    time.sleep(delay)
                    continue
                return q
            _mark_api("http=200")
            choices = result.get('choices') or []
            if not choices:
                retries += 1
                if retries < max_retries:
                    delay = min(retry_base_delay * (2 ** (retries - 1)), retry_max_delay)
                    delay = delay + (random.random() * 0.5)
                    time.sleep(delay)
                    continue
                return q
            message = (choices[0] or {}).get('message') or {}
            content = message.get('content')
            if content is None:
                retries += 1
                if retries < max_retries:
                    delay = min(retry_base_delay * (2 ** (retries - 1)), retry_max_delay)
                    delay = delay + (random.random() * 0.5)
                    time.sleep(delay)
                    continue
                return q
            translated = str(content).strip()
            if strict_no_hangul and (translated == original_stripped or HANGUL_RE.search(translated)):
                retries += 1
                if retries < max_retries:
                    user_prompt = (
                        f"Translate from {from_lang} to {to_lang}.\n"
                        "Return ONLY the translated text.\n"
                        "Do NOT leave any Korean characters in the output.\n"
                        "Text:\n"
                        f"{q}"
                    )
                    payload['messages'] = [
                        {'role': 'system', 'content': system_prompt},
                        {'role': 'user', 'content': user_prompt},
                    ]
                    delay = min(retry_base_delay * (2 ** (retries - 1)), retry_max_delay)
                    delay = delay + (random.random() * 0.5)
                    time.sleep(delay)
                    continue
                return q
            _translation_cache[cache_key] = translated
            return translated
        except Exception:
            _mark_api("exception")
            retries += 1
            if retries < max_retries:
                delay = min(retry_base_delay * (2 ** (retries - 1)), retry_max_delay)
                delay = delay + (random.random() * 0.5)
                time.sleep(delay)
    return q

def get_translation(text):
    if not text or not text.strip(): return text
    direction = translation_direction.get()
    from_lang = FROM_LANG_MAP.get(direction, 'auto')
    to_lang = TO_LANG_MAP.get(direction, 'auto')
    
    translated_text = kimi_translate(text, from_lang, to_lang)
    # 事后校准
    translated_text = apply_revisions(translated_text)
    return translated_text

def _ppt_normalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")

def _ppt_denormalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\n", "\v")

def translate_ppt_paragraph(paragraph):
    full_text = paragraph.text
    if not full_text or not str(full_text).strip():
        return

    normalized = _ppt_normalize_linebreaks(full_text)
    translated = get_translation(normalized)

    original_texts.append(normalized)
    translated_texts.append(translated)

    if append_translation.get():
        result_norm = append_translation_to_original(normalized, translated)
    else:
        result_norm = translated

    result_text = _ppt_denormalize_linebreaks(result_norm)

    if paragraph.runs:
        first_run = paragraph.runs[0]
    else:
        first_run = paragraph.add_run()

    original_font_name = first_run.font.name
    original_size = first_run.font.size
    original_bold = first_run.font.bold
    original_italic = first_run.font.italic
    original_underline = first_run.font.underline

    first_run.text = result_text
    for r in paragraph.runs[1:]:
        r.text = ""

    target_font_name = get_target_font_name()
    if target_font_name:
        set_pptx_run_font(first_run, target_font_name)
        set_pptx_paragraph_default_font(paragraph, target_font_name)
    elif original_font_name:
        first_run.font.name = original_font_name
        first_run.font.size = original_size
        first_run.font.bold = original_bold
        first_run.font.italic = original_italic
        first_run.font.underline = original_underline

def append_translation_to_original(text, translated_text, cell=None):
    text = text.strip()
    translated_text = translated_text.strip()
    result = f"{text}\n{translated_text}" if text and translated_text else (text or translated_text)
    if cell:
        cell.alignment = Alignment(wrap_text=True, vertical='center')
        ws = cell.parent
        row_num = cell.row
        original_height = ws.row_dimensions[row_num].height
        ws.row_dimensions[row_num].height = (original_height * 2) if original_height else 30
    return result

# ==========================================
# 文件处理逻辑 (完全基于 V2.9 的稳定代码)
# ==========================================

def translate_shape_for_ppt(shape):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            translate_ppt_paragraph(paragraph)
    elif shape.shape_type == 6:  # 组合形状
        for sub_shape in shape.shapes:
            translate_shape_for_ppt(sub_shape)
    elif shape.has_table:  # 处理表格形状
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    translate_ppt_paragraph(paragraph)

def update_ui_status(msg):
    """线程安全地更新 UI 状态"""
    _mark_progress(msg)
    root.after(0, lambda: status_label.config(text=msg))

def translate_ppt(input_file, output_file):
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        msg = f"正在翻译 PPT: 第 {i}/{total_slides} 页..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        for shape in slide.shapes:
            translate_shape_for_ppt(shape)
    
    print(f"[系统] 正在保存 PPT 文件...")
    prs.save(output_file)
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] 保存修改后的 PPT 文件完成!")

def clean_sheet_name(name):
    if not name: return "Sheet"
    invalid_chars = r'[\\/?:*\[\](){}<>|"\']'
    clean_name = re.sub(invalid_chars, '', name)
    return clean_name[:31]

def translate_excel_xlsx(input_file, output_file):
    """处理.xlsx格式的Excel文件"""
    # 加载工作簿 (keep_vba=True 尝试保留宏和部分绘图元数据)
    wb = load_workbook(input_file, keep_vba=True)
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    total_sheets = len(wb.sheetnames)
    for i, sheet_name in enumerate(wb.sheetnames, 1):
        msg = f"正在翻译 Excel: 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        # 翻译工作表名称 (还原 V2.9 逻辑并增加翻译)
        translated_sheet_name = get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        ws = wb[sheet_name]
        # 修改工作表名称
        ws.title = unique_sheet_name
        
        # 获取总行数用于进度（大致）
        total_rows = ws.max_row
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            if row_idx % 10 == 0:
                row_msg = f"处理中: {unique_sheet_name} 第 {row_idx}/{total_rows} 行..."
                print(f"[进度] {row_msg}")
                update_ui_status(row_msg)
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    original_value = str(cell.value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    if "\n" in normalized_value:
                        translated_lines = []
                        for line in normalized_value.split("\n"):
                            if line.strip():
                                translated_lines.append(get_translation(line))
                            else:
                                translated_lines.append(line)
                        translated_text = "\n".join(translated_lines)
                    else:
                        translated_text = get_translation(normalized_value)

                    original_texts.append(normalized_value)
                    translated_texts.append(translated_text)
                    
                    # 严格还原 V2.9 样式保护逻辑
                    # 保存原格式
                    original_font = Font(**cell.font.__dict__)
                    original_border = Border(**cell.border.__dict__)
                    original_alignment = Alignment(**cell.alignment.__dict__)
                    original_fill = PatternFill(**cell.fill.__dict__)
                    
                    if append_translation.get():
                        cell.value = append_translation_to_original(normalized_value, translated_text, cell)
                    else:
                        cell.value = translated_text
                    
                    # 恢复原格式
                    target_font_name = get_target_font_name()
                    if target_font_name:
                        try:
                            cell.font = original_font.copy(name=target_font_name)
                        except Exception:
                            cell.font = Font(
                                name=target_font_name,
                                size=original_font.size,
                                bold=original_font.bold,
                                italic=original_font.italic,
                                underline=original_font.underline,
                                color=original_font.color,
                            )
                    else:
                        cell.font = original_font
                    cell.border = original_border
                    cell.alignment = original_alignment
                    cell.fill = original_fill
                    # 设置单元格自动换行
                    cell.alignment = Alignment(wrap_text=True, vertical='center')
    # 保存修改后的工作簿
    wb.save(output_file)
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)

def translate_excel_xls(input_file, output_file):
    """处理.xls格式的Excel文件，使用pandas库，保存为.xlsx格式"""
    base_output = os.path.splitext(output_file)[0]
    output_file_xlsx = base_output + '.xlsx'
    target_font_name = get_target_font_name()
    # 读取所有工作表
    excel_file = pd.ExcelFile(input_file)
    # 创建一个新的ExcelWriter对象，使用默认引擎
    writer = pd.ExcelWriter(output_file_xlsx, engine='openpyxl')
    # 收集所有已存在的工作表名称（包括原始和已处理的）
    existing_sheet_names = set()
    total_sheets = len(excel_file.sheet_names)
    for i, sheet_name in enumerate(excel_file.sheet_names, 1):
        msg = f"正在翻译 Excel(.xls): 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        # 翻译工作表名称
        translated_sheet_name = get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        # 确保工作表名称唯一
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        # 读取工作表
        df = pd.read_excel(input_file, sheet_name=sheet_name)
        # 遍历所有单元格进行翻译
        total_rows = len(df.index)
        for idx_num, idx in enumerate(df.index, 1):
            if idx_num % 10 == 0:
                row_msg = f"处理中: {unique_sheet_name} 第 {idx_num}/{total_rows} 行..."
                print(f"[进度] {row_msg}")
                update_ui_status(row_msg)
            for col in df.columns:
                cell_value = df.at[idx, col]
                if pd.notna(cell_value) and isinstance(cell_value, str):
                    original_value = str(cell_value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    if "\n" in normalized_value:
                        translated_lines = []
                        for line in normalized_value.split("\n"):
                            if line.strip():
                                translated_lines.append(get_translation(line))
                            else:
                                translated_lines.append(line)
                        translated_text = "\n".join(translated_lines)
                    else:
                        translated_text = get_translation(normalized_value)

                    original_texts.append(normalized_value)
                    translated_texts.append(translated_text)

                    df.at[idx, col] = append_translation_to_original(normalized_value, translated_text) if append_translation.get() else translated_text
        # 保存翻译后的工作表
        df.to_excel(writer, sheet_name=unique_sheet_name, index=False)
    # 保存工作簿
    writer.close()
    if target_font_name:
        try:
            wb = load_workbook(output_file_xlsx)
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            try:
                                cell.font = cell.font.copy(name=target_font_name)
                            except Exception:
                                cell.font = Font(
                                    name=target_font_name,
                                    size=cell.font.size,
                                    bold=cell.font.bold,
                                    italic=cell.font.italic,
                                    underline=cell.font.underline,
                                    color=cell.font.color,
                                )
            wb.save(output_file_xlsx)
        except Exception:
            pass
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)
    return output_file_xlsx

def translate_word(input_file, output_file):
    doc = Document(input_file)
    target_font_name = get_target_font_name()

    def set_style_font(style_element, font_name):
        if style_element is None or not font_name:
            return
        rPr = style_element.find(qn('w:rPr'))
        if rPr is None:
            rPr = OxmlElement('w:rPr')
            style_element.append(rPr)
        rFonts = rPr.find(qn('w:rFonts'))
        if rFonts is None:
            rFonts = OxmlElement('w:rFonts')
            rPr.append(rFonts)
        rFonts.set(qn('w:ascii'), font_name)
        rFonts.set(qn('w:hAnsi'), font_name)
        rFonts.set(qn('w:eastAsia'), font_name)
        rFonts.set(qn('w:cs'), font_name)
        for k in (qn('w:asciiTheme'), qn('w:hAnsiTheme'), qn('w:eastAsiaTheme'), qn('w:csTheme')):
            try:
                rFonts.attrib.pop(k, None)
            except Exception:
                pass

    def translate_word_paragraph(paragraph):
        full_text = paragraph.text
        if not full_text or not str(full_text).strip():
            return

        normalized = str(full_text).replace("\r\n", "\n").replace("\r", "\n")
        t = get_translation(normalized)

        original_texts.append(normalized)
        translated_texts.append(t)

        result_text = append_translation_to_original(normalized, t) if append_translation.get() else t

        drawings = []
        for run in paragraph.runs:
            for inline in xpath_with_ns(run.element, './/w:drawing'):
                pic_elements = xpath_with_ns(inline, './/a:blip/@r:embed')
                if pic_elements:
                    drawings.append(inline)

        if paragraph.runs:
            first_run = paragraph.runs[0]
        else:
            first_run = paragraph.add_run()

        first_run.text = result_text
        if target_font_name:
            set_docx_run_font(first_run, target_font_name)

        for run in paragraph.runs[1:]:
            has_drawing = bool(xpath_with_ns(run.element, './/w:drawing'))
            if not has_drawing:
                run.text = ""

        for inline in drawings:
            new_run = paragraph.add_run()
            new_run._r.append(inline)

    def translate_word_table(table):
        raise RuntimeError("translate_word_table should not be called directly")

    def count_table_cells(table):
        total = 0
        for row in table.rows:
            for cell in row.cells:
                total += 1
                for nested_table in cell.tables:
                    total += count_table_cells(nested_table)
        return total

    def translate_word_table_with_progress(table, table_idx, total_tables, total_cells, cell_counter, progress_state):
        now = time.time()
        if now - progress_state[0] >= 30:
            msg = f"正在翻译 Word: 表格 {table_idx}/{total_tables} 单元格 {cell_counter[0]}/{total_cells}..."
            print(f"[进度] {msg}")
            update_ui_status(msg)
            progress_state[0] = now
        for row in table.rows:
            for cell in row.cells:
                cell_counter[0] += 1
                if cell_counter[0] % 20 == 0:
                    msg = f"正在翻译 Word: 表格 {table_idx}/{total_tables} 单元格 {cell_counter[0]}/{total_cells}..."
                    print(f"[进度] {msg}")
                    update_ui_status(msg)
                    progress_state[0] = time.time()
                for paragraph in cell.paragraphs:
                    translate_word_paragraph(paragraph)
                for nested_table in cell.tables:
                    translate_word_table_with_progress(nested_table, table_idx, total_tables, total_cells, cell_counter, progress_state)
    
    # 处理普通段落 (完全还原 V2.9 逻辑)
    total_paragraphs = len(doc.paragraphs)
    for i, paragraph in enumerate(doc.paragraphs, 1):
        if i % 5 == 0:
            msg = f"正在翻译 Word: 段落 {i}/{total_paragraphs}..."
            print(f"[进度] {msg}")
            update_ui_status(msg)
        translate_word_paragraph(paragraph)
    
    # 处理表格 (完全还原 V2.9 逻辑)
    total_tables = len(doc.tables)
    for i, table in enumerate(doc.tables, 1):
        msg = f"正在翻译 Word: 表格 {i}/{total_tables}..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        total_cells = count_table_cells(table)
        cell_counter = [0]
        progress_state = [time.time()]
        translate_word_table_with_progress(table, i, total_tables, total_cells, cell_counter, progress_state)
    
    # 处理形状中的文字 (完全还原 V2.9 逻辑)
    for shape in doc.inline_shapes:
        try:
            # 通过 XML 访问形状中的文本
            if hasattr(shape, '_inline') and hasattr(shape._inline, 'graphic'):
                graphic_data = shape._inline.graphic.graphicData
                if hasattr(graphic_data, 'txBody') and graphic_data.txBody:
                    for p in graphic_data.txBody.p:
                        for r in p.r:
                            if hasattr(r, 't'):
                                text = r.t
                                if text and text.strip():  # 只翻译非空文本
                                    # 先获取翻译结果
                                    t = get_translation(text)
                                    
                                    # 收集翻译前后的文本
                                    original_texts.append(text)
                                    translated_texts.append(t)
                                    if append_translation.get():
                                        r.t = append_translation_to_original(text, t)
                                    else:
                                        r.t = t
                                    if target_font_name:
                                        set_drawingml_r_element_font(r, target_font_name)
        except Exception:
            continue
    doc_element_list = [doc.element]
    for section in doc.sections:
        doc_element_list.extend([
            section.header._element,
            section.footer._element,
            section.first_page_header._element,
            section.first_page_footer._element,
            section.even_page_header._element,
            section.even_page_footer._element,
        ])
    for doc_element in doc_element_list:
        for t_node in xpath_with_ns(doc_element, './/w:txbxContent//w:t'):
            text = t_node.text
            if not text or not text.strip():
                continue
            t = get_translation(text)
            original_texts.append(text)
            translated_texts.append(t)
            t_node.text = append_translation_to_original(text, t) if append_translation.get() else t
            if not target_font_name:
                continue
            r_element = t_node
            while r_element is not None and r_element.tag != qn('w:r'):
                r_element = r_element.getparent()
            if r_element is not None:
                set_docx_r_element_font(r_element, target_font_name)
        if not append_translation.get():
            for txbx in xpath_with_ns(doc_element, './/w:txbxContent'):
                t_nodes = xpath_with_ns(txbx, './/w:t')
                if not t_nodes:
                    continue
                combined = ''.join([(n.text or '') for n in t_nodes])
                revised = apply_revisions(combined)
                if revised != combined:
                    t_nodes[0].text = revised
                    for n in t_nodes[1:]:
                        n.text = ''
                    if target_font_name:
                        r_element = t_nodes[0]
                        while r_element is not None and r_element.tag != qn('w:r'):
                            r_element = r_element.getparent()
                        if r_element is not None:
                            set_docx_r_element_font(r_element, target_font_name)
        for t_node in doc_element.xpath(f'.//*[namespace-uri()="{A_NS}" and local-name()="txBody"]//*[namespace-uri()="{A_NS}" and local-name()="t"]'):
            text = t_node.text
            if not text or not text.strip():
                continue
            t = get_translation(text)
            original_texts.append(text)
            translated_texts.append(t)
            t_node.text = append_translation_to_original(text, t) if append_translation.get() else t
            if not target_font_name:
                continue
            a_r_element = t_node
            while a_r_element is not None and a_r_element.tag != f'{{{A_NS}}}r':
                a_r_element = a_r_element.getparent()
            if a_r_element is not None:
                set_drawingml_r_element_font(a_r_element, target_font_name)
        if not append_translation.get():
            for a_txbody in doc_element.xpath(f'.//*[namespace-uri()="{A_NS}" and local-name()="txBody"]'):
                t_nodes = a_txbody.xpath(f'.//*[namespace-uri()="{A_NS}" and local-name()="t"]')
                if not t_nodes:
                    continue
                combined = ''.join([(n.text or '') for n in t_nodes])
                revised = apply_revisions(combined)
                if revised != combined:
                    t_nodes[0].text = revised
                    for n in t_nodes[1:]:
                        n.text = ''
                    if target_font_name:
                        a_r_element = t_nodes[0]
                        while a_r_element is not None and a_r_element.tag != f'{{{A_NS}}}r':
                            a_r_element = a_r_element.getparent()
                        if a_r_element is not None:
                            set_drawingml_r_element_font(a_r_element, target_font_name)
        if target_font_name:
            for r_element in xpath_with_ns(doc_element, './/w:txbxContent//w:r'):
                set_docx_r_element_font(r_element, target_font_name)
            for r_element in xpath_with_ns(doc_element, './/w:r'):
                set_docx_r_element_font(r_element, target_font_name)
            for a_rpr_element in doc_element.xpath(f'.//*[namespace-uri()="{A_NS}" and (local-name()="rPr" or local-name()="defRPr" or local-name()="endParaRPr")]'):
                set_drawingml_rpr_element_font(a_rpr_element, target_font_name)
            try:
                for s in doc.styles:
                    try:
                        if getattr(s, "font", None) is not None:
                            s.font.name = target_font_name
                    except Exception:
                        pass
                    try:
                        set_style_font(getattr(s, "_element", None), target_font_name)
                    except Exception:
                        pass
            except Exception:
                pass

    doc.save(output_file)
    # 翻译完成后调用保存语料库函数
    save_to_corpus(original_texts, translated_texts)

# --- 线程控制 ---

def start_translation():
    input_file = input_file_entry.get()
    output_folder = output_folder_entry.get()
    if not (input_file and output_folder):
        messagebox.showwarning("提示", "请完整选择输入文件和输出目录")
        return
    translate_button.config(state=tk.DISABLED, text="🚀 正在翻译，请稍候...")
    status_label.config(text="任务已启动，请查看终端进度...", foreground="#2980b9")
    thread = threading.Thread(target=run_translation_task, args=(input_file, output_folder))
    thread.daemon = True
    thread.start()

def run_translation_task(input_file, output_folder):
    _watchdog_start()
    try:
        global revision_map
        file_ext = os.path.splitext(input_file)[1].lower()
        custom_name = custom_filename_entry.get().strip()
        output_file = os.path.join(output_folder, (custom_name if custom_name else f"translated_v2.13_{os.path.basename(input_file).split('.')[0]}") + file_ext)
        
        # --- [V2.13 增强功能：物理克隆以保留图片和绘图] ---
        # 即使 V2.9 也不支持 Excel 绘图保留，这里通过物理复制尝试最大化兼容性
        try:
            if os.path.exists(output_file):
                os.remove(output_file) # 尝试删除旧文件，如果被占用会在这里报错
            shutil.copy2(input_file, output_file)
        except PermissionError:
            raise Exception(f"目标文件已被占用，请先关闭 Excel/Word/PPT: {os.path.basename(output_file)}")
        
        original_texts.clear()
        translated_texts.clear()
        revision_map = load_revision_dict("revision.md", silent=True)
        print(f"\n[系统] 开始翻译任务: {os.path.basename(input_file)}")
        print(f"[系统] 校准规则加载: {len(revision_map)} 条")
        update_ui_status(f"开始翻译: {os.path.basename(input_file)}")
        
        # 统一使用 output_file 作为操作对象，实现“原地翻译”
        if file_ext in ['.ppt', '.pptx']: translate_ppt(output_file, output_file)
        elif file_ext == '.xlsx': translate_excel_xlsx(output_file, output_file)
        elif file_ext == '.xls': 
            # .xls 比较特殊，必须另存为 .xlsx
            output_file = translate_excel_xls(input_file, output_file)
        elif file_ext == '.docx': translate_word(output_file, output_file)
        
        print(f"[完成] 文件已保存至: {output_file}\n")
        root.after(0, lambda: translation_done_callback(output_file))
    except Exception as e:
        err_msg = str(e)
        print(f"[错误] 详情: {err_msg}")
        root.after(0, lambda: translation_failed_callback(err_msg))
    finally:
        _watchdog_stop()

def translation_done_callback(output_file):
    translate_button.config(state=tk.NORMAL, text="🚀 开始 KIMI 翻译")
    status_label.config(text=f"翻译任务已圆满完成！", foreground="#27ae60")
    messagebox.showinfo("成功", f"翻译完成！\n文件保存至：{output_file}")

def translation_failed_callback(error_msg):
    translate_button.config(state=tk.NORMAL, text="🚀 开始 KIMI 翻译")
    status_label.config(text=f"翻译过程出错", foreground="#e74c3c")
    messagebox.showerror("错误", f"发生异常：{error_msg}")

def save_to_corpus(orig, trans):
    if generate_corpus.get() and orig:
        if not os.path.exists('Corpus'): os.makedirs('Corpus')
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        try:
            direction = translation_direction.get()
        except Exception:
            direction = "unknown"
        corpus_file = f'Corpus/Corpus_v2.13_{direction}_{timestamp}.xlsx'
        to_lang = TO_LANG_MAP.get(direction)
        seen = set()
        filtered_orig = []
        filtered_trans = []
        for o, t in zip(orig, trans):
            if o is None or t is None:
                continue
            o = str(o).replace("\r\n", "\n").replace("\r", "\n").strip()
            t = str(t).replace("\r\n", "\n").replace("\r", "\n").strip()
            if not o or not t:
                continue
            if to_lang in ("zh", "zh_tw") and re.search(r"[\u4e00-\u9fff]", o):
                continue
            if o == t:
                continue
            if not re.search(r"[\uac00-\ud7a3]", o) and re.search(r"[A-Za-z0-9]", o):
                continue
            if not re.search(r"[\uac00-\ud7a3\u4e00-\u9fff]", o) and not re.search(r"[A-Za-z0-9]", o):
                continue
            key = o
            if key in seen:
                continue
            seen.add(key)
            filtered_orig.append(o)
            filtered_trans.append(t)

        if not filtered_orig:
            return

        pd.DataFrame(
            {'序号': range(1, len(filtered_orig) + 1), '翻译前': filtered_orig, '翻译后': filtered_trans}
        ).to_excel(corpus_file, index=False)

# --- UI 布局 ---
root = tk.Tk()
root.title("KIMI 翻译工具 V2.13")
root.geometry("700x850")
root.configure(bg="#f5f6fa")
style = ttk.Style()
style.theme_use('clam')
style.configure("TFrame", background="#f5f6fa")
style.configure("TLabel", background="#f5f6fa", font=("微软雅黑", 10))
style.configure("Header.TLabel", font=("微软雅黑", 14, "bold"), foreground="#166534")
main_frame = ttk.Frame(root, padding="20")
main_frame.pack(fill="both", expand=True)
ttk.Label(main_frame, text="KIMI 翻译 & 自动校准系统", style="Header.TLabel").pack(pady=(0, 20))
file_card = ttk.LabelFrame(main_frame, text=" 文件设置 ", padding=15)
file_card.pack(fill="x", pady=10)
ttk.Label(file_card, text="待翻译文件:").grid(row=0, column=0, sticky="w", pady=5)
input_file_entry = ttk.Entry(file_card, width=50)
input_file_entry.grid(row=0, column=1, padx=10)
ttk.Button(file_card, text="选择文件", command=lambda: (input_file_entry.delete(0, tk.END), input_file_entry.insert(0, filedialog.askopenfilename()))).grid(row=0, column=2)
ttk.Label(file_card, text="保存位置:").grid(row=1, column=0, sticky="w", pady=5)
output_folder_entry = ttk.Entry(file_card, width=50)
output_folder_entry.grid(row=1, column=1, padx=10)
ttk.Button(file_card, text="选择目录", command=lambda: (output_folder_entry.delete(0, tk.END), output_folder_entry.insert(0, filedialog.askdirectory()))).grid(row=1, column=2)
ttk.Label(file_card, text="自定义文件名:").grid(row=2, column=0, sticky="w", pady=5)
custom_filename_entry = ttk.Entry(file_card, width=50)
custom_filename_entry.grid(row=2, column=1, padx=10, columnspan=2, sticky="w")
dir_card = ttk.LabelFrame(main_frame, text=" 翻译语种 ", padding=15)
dir_card.pack(fill="x", pady=10)
translation_direction = tk.StringVar(value='ko2zh')
lang_grid = ttk.Frame(dir_card)
lang_grid.pack(fill="x")
langs = [("韩 -> 中", 'ko2zh'), ("中 -> 韩", 'zh2ko'), ("韩 -> 越", 'ko2vi'), ("韩 -> 英", 'ko2en'),
         ("中 -> 英", 'zh2en'), ("英 -> 中", 'en2zh'), ("繁中 -> 英", 'zh_tw2en'), ("英 -> 繁中", 'en2zh_tw'),
         ("中 -> 日", 'zh2ja'), ("日 -> 中", 'ja2zh'), ("英 -> 韩", 'en2ko'), ("越 -> 中", 'vi2zh')]
for i, (text, val) in enumerate(langs):
    row, col = i // 4, i % 4
    ttk.Radiobutton(lang_grid, text=text, variable=translation_direction, value=val).grid(row=row, column=col, padx=15, pady=5, sticky="w")
opt_frame = tk.Frame(main_frame, bg="#f5f6fa")
opt_frame.pack(fill="x", pady=10)
append_translation = tk.BooleanVar(value=False)
tk.Checkbutton(opt_frame, text="在原文下方保留翻译对照", variable=append_translation, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
generate_corpus = tk.BooleanVar(value=False)
tk.Checkbutton(opt_frame, text="同步生成语料库 (Corpus)", variable=generate_corpus, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
translate_button = tk.Button(main_frame, text="🚀 开始 KIMI 翻译", command=start_translation, bg="#16a34a", fg="white", font=("微软雅黑", 12, "bold"), relief="flat", cursor="hand2", pady=12)
translate_button.pack(fill="x", pady=20)
status_label = ttk.Label(main_frame, text="就绪：已加载校准词典 (revision.md)", foreground="#7f8c8d")
status_label.pack()

# 初始化加载一次即可
revision_map = load_revision_dict("revision.md")

root.mainloop()
