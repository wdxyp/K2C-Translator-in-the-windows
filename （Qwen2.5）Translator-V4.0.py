import warnings
# 静默 openpyxl 关于 DrawingML 支持不全的警告
warnings.filterwarnings("ignore", category=UserWarning, module="openpyxl.reader.drawings")

import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
import pandas as pd
import os
import json
import hashlib
import random
import string
import requests
from openpyxl import load_workbook
from openpyxl.styles import Font, Border, Alignment, PatternFill
from docx import Document
from docx.enum.text import WD_BREAK
from docx.text.paragraph import Paragraph
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.oxml.ns import nsmap as DOCX_NSMAP
from datetime import datetime
import re
import time
import threading
import shutil
import copy
import traceback
import urllib.request
from concurrent.futures import ThreadPoolExecutor, as_completed

# ==========================================
# V4.0 (Qwen2.5 本地模型翻译版)
# ==========================================

OLLAMA_BASE_URL = (os.getenv("OLLAMA_BASE_URL") or "http://127.0.0.1:11434").strip().rstrip("/")
OLLAMA_MODEL_NAME = (os.getenv("OLLAMA_MODEL") or "qwen2.5:3b").strip()
QWEN_DEFAULT_BATCH_SIZE = int(os.getenv("QWEN_BATCH_SIZE", "8"))
QWEN_DEFAULT_MAX_NEW_TOKENS = int(os.getenv("QWEN_MAX_NEW_TOKENS", "128"))
QWEN_DEFAULT_PROMPT_MAX_LENGTH = int(os.getenv("QWEN_PROMPT_MAX_LENGTH", "2048"))
OLLAMA_TIMEOUT = int(os.getenv("OLLAMA_TIMEOUT", "180"))
OLLAMA_MAX_WORKERS = int(os.getenv("OLLAMA_MAX_WORKERS", "2"))

# 全局变量
revision_map = {}
original_texts = []
translated_texts = []
# 扩大韩文匹配范围：包括预组合音节、基础辅音/元音、扩展字母等
HANGUL_RE = re.compile(r"[\uac00-\ud7af\u1100-\u11ff\u3130-\u318f\ua960-\ua97f\ud7b0-\ud7ff]")
# 扩大中文匹配范围：包括基本、扩展A、兼容汉字等
CHINESE_RE = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf\uf900-\ufaff]")
_ollama_runtime_lock = threading.Lock()
_translation_cache = {}  # 内存缓存：记录已翻译过的文本，避免重复请求 API
_qwen_batch_size = QWEN_DEFAULT_BATCH_SIZE
_kimi_max_workers = 1  # 保留原有文件处理流程接口，但本地模型默认单线程更稳定

LANGUAGE_NAME_MAP = {
    "kor": "Korean",
    "zh": "Simplified Chinese",
    "zh_tw": "Traditional Chinese",
    "en": "English",
    "ja": "Japanese",
    "vie": "Vietnamese",
}

# #region debug-point A:runtime-hooks
def _debug_report(hypothesis_id, location, msg, data=None):
    try:
        env_path = os.path.join(os.path.dirname(__file__), ".dbg", "qwen-auto-exit.env")
        url = "http://127.0.0.1:7777/event"
        session_id = "qwen-auto-exit"
        try:
            with open(env_path, "r", encoding="utf-8") as f:
                content = f.read()
            for line in content.splitlines():
                if line.startswith("DEBUG_SERVER_URL="):
                    url = line.split("=", 1)[1].strip() or url
                elif line.startswith("DEBUG_SESSION_ID="):
                    session_id = line.split("=", 1)[1].strip() or session_id
        except Exception:
            pass
        payload = {
            "sessionId": session_id,
            "runId": "pre-fix",
            "hypothesisId": hypothesis_id,
            "location": location,
            "msg": f"[DEBUG] {msg}",
            "data": data or {},
            "ts": int(time.time() * 1000),
        }
        req = urllib.request.Request(
            url,
            data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
            headers={"Content-Type": "application/json"},
        )
        urllib.request.urlopen(req, timeout=1).read()
    except Exception:
        pass


def _debug_threading_excepthook(args):
    _debug_report(
        "B",
        "threading.excepthook",
        "线程未捕获异常",
        {
            "thread_name": getattr(args.thread, "name", ""),
            "exc_type": getattr(getattr(args, "exc_type", None), "__name__", ""),
            "exc_value": str(getattr(args, "exc_value", "")),
            "traceback": "".join(traceback.format_exception(args.exc_type, args.exc_value, args.exc_traceback))[:4000],
        },
    )


threading.excepthook = _debug_threading_excepthook
_debug_report("A", "module-import", "模块已加载", {"pid": os.getpid(), "python": os.sys.version})
# #endregion

def should_skip_translation(text, to_lang):
    """
    智能判断是否跳过翻译：
    1. 如果目标是中文且原文已含中文，则跳过（满足用户：即使含部分韩文也不翻译中文）
    2. 如果原文和目标语种特征一致，则跳过
    """
    if not text or not str(text).strip():
        return True
    
    # 转换为字符串处理
    text_s = str(text).strip()
    
    # 核心逻辑：如果目标是中文（简体/繁体），且原文中已经包含中文字符，则跳过翻译直接返回原文
    if to_lang in ('zh', 'zh_tw') and CHINESE_RE.search(text_s):
        return True
        
    # 如果目标是韩文，且原文全是韩文（不含中文），也可以考虑跳过
    if to_lang == 'kor' and HANGUL_RE.search(text_s) and not CHINESE_RE.search(text_s):
        return True
        
    # 如果目标是英文，且原文不含中韩文，则跳过
    if to_lang == 'en' and not HANGUL_RE.search(text_s) and not CHINESE_RE.search(text_s):
        return True
        
    return False

def _normalize_ollama_base_url(base_url: str) -> str:
    normalized = str(base_url or "").strip().strip('"').strip("'").rstrip("/")
    if not normalized:
        raise RuntimeError("Ollama 服务地址不能为空")
    if not normalized.lower().startswith(("http://", "https://")):
        normalized = "http://" + normalized
    return normalized.rstrip("/")


def _check_ollama_runtime():
    with _ollama_runtime_lock:
        base_url = _normalize_ollama_base_url(OLLAMA_BASE_URL)
        model_name = str(OLLAMA_MODEL_NAME or "").strip()
        if not model_name:
            raise RuntimeError("Ollama 模型名不能为空")

        # #region debug-point C:ollama-runtime-check
        _debug_report("C", "_check_ollama_runtime", "检查 Ollama 运行状态", {"base_url": base_url, "model": model_name})
        # #endregion

        try:
            resp = requests.get(f"{base_url}/api/tags", timeout=10)
            resp.raise_for_status()
            tags_payload = resp.json()
        except Exception as e:
            raise RuntimeError(f"Ollama 服务不可用: {base_url}，请先启动 Ollama。错误: {e}")

        models = tags_payload.get("models") or []
        model_names = {str(item.get("name", "")).strip() for item in models if isinstance(item, dict)}
        if model_name not in model_names:
            raise RuntimeError(
                f"Ollama 模型不存在: {model_name}。"
                f" 请先执行 `ollama pull {model_name}`，或在界面中改成已安装模型名。"
            )
        return base_url, model_name

_last_progress_ts = 0.0
_last_progress_msg = ""
_last_api_ts = 0.0
_last_api_status = ""
_watchdog_stop_event = threading.Event()
_watchdog_thread = None
def _mark_progress(msg):
    global _last_progress_ts, _last_progress_msg
    _last_progress_ts = time.time()
    _last_progress_msg = str(msg)

def _mark_api(status):
    global _last_api_ts, _last_api_status
    _last_api_ts = time.time()
    _last_api_status = str(status)

def _watchdog_loop():
    while not _watchdog_stop_event.is_set():
        time.sleep(20)
        now = time.time()
        last_progress = _last_progress_ts
        last_api = _last_api_ts
        msg = _last_progress_msg or "运行中"
        progress_age = int(now - last_progress) if last_progress else -1
        api_age = int(now - last_api) if last_api else -1
        status = _last_api_status or "unknown"
        if progress_age >= 60 or api_age >= 60:
            print(f"[心跳] {msg} | 距上次进度 {progress_age}s | 距上次翻译 {api_age}s | 上次状态 {status}")

def _watchdog_start():
    global _watchdog_thread
    _watchdog_stop_event.clear()
    _watchdog_thread = threading.Thread(target=_watchdog_loop, daemon=True)
    _watchdog_thread.start()

def _watchdog_stop():
    _watchdog_stop_event.set()

FROM_LANG_MAP = {
    'zh2ko': 'zh', 'ko2zh': 'kor', 'ko2en': 'kor', 'zh2en': 'zh', 'en2zh': 'en',
    'zh_tw2en': 'zh', 'en2zh_tw': 'en', 'zh2ja': 'zh', 'ja2zh': 'ja',
    'en2ko': 'en', 'vi2zh': 'vie', 'zh2vi': 'zh', 'ko2vi': 'kor'
}
TO_LANG_MAP = {
    'zh2ko': 'kor', 'ko2zh': 'zh', 'ko2en': 'en', 'zh2en': 'en', 'en2zh': 'zh',
    'zh_tw2en': 'en', 'en2zh_tw': 'zh', 'zh2ja': 'ja', 'ja2zh': 'zh',
    'en2ko': 'kor', 'vi2zh': 'zh', 'zh2vi': 'vie', 'ko2vi': 'vie'
}
TARGET_FONT_BY_TO_LANG = {'zh': '微软雅黑', 'kor': 'Malgun Gothic'}
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def get_target_to_lang():
    try:
        if root and root.winfo_exists():
            direction = translation_direction.get()
            return TO_LANG_MAP.get(direction)
    except Exception:
        pass
    return None

def get_target_font_name():
    to_lang = get_target_to_lang()
    if not to_lang:
        return None
    return TARGET_FONT_BY_TO_LANG.get(to_lang)

def xpath_with_ns(element, expr):
    if element is None:
        return []
    try:
        return element.xpath(expr, namespaces=DOCX_NSMAP)
    except TypeError:
        return element.xpath(expr)

def set_docx_r_element_font(r_element, font_name):
    if r_element is None or not font_name:
        return
    rPr = r_element.get_or_add_rPr()
    rFonts = rPr.rFonts
    if rFonts is None:
        rFonts = OxmlElement('w:rFonts')
        rPr.append(rFonts)
    rFonts.set(qn('w:ascii'), font_name)
    rFonts.set(qn('w:hAnsi'), font_name)
    rFonts.set(qn('w:eastAsia'), font_name)
    rFonts.set(qn('w:cs'), font_name)
    for k in (qn('w:asciiTheme'), qn('w:hAnsiTheme'), qn('w:eastAsiaTheme'), qn('w:csTheme')):
        try:
            rFonts.attrib.pop(k, None)
        except Exception:
            pass

def set_docx_run_font(run, font_name):
    if not run or not font_name:
        return
    run.font.name = font_name
    set_docx_r_element_font(run._r, font_name)

def apply_run_format(source_run, target_run):
    """将 source_run 的所有格式（rPr）深度应用到 target_run"""
    if not source_run or not target_run:
        return
    source_rPr = source_run._element.rPr
    if source_rPr is not None:
        new_rPr = copy.deepcopy(source_rPr)
        target_r = target_run._element
        if target_r.rPr is not None:
            target_r.remove(target_r.rPr)
        target_r.insert(0, new_rPr)

def set_drawingml_rpr_element_font(a_rPr_element, font_name):
    if a_rPr_element is None or not font_name:
        return
    for local_name, tag in (('latin', 'a:latin'), ('ea', 'a:ea'), ('cs', 'a:cs')):
        el = a_rPr_element.find(f'{{{A_NS}}}{local_name}')
        if el is None:
            el = OxmlElement(tag)
            a_rPr_element.append(el)
        el.set('typeface', font_name)

def set_drawingml_r_element_font(a_r_element, font_name):
    if a_r_element is None or not font_name:
        return
    a_rPr = a_r_element.find(f'{{{A_NS}}}rPr')
    if a_rPr is None:
        a_rPr = OxmlElement('a:rPr')
        a_r_element.insert(0, a_rPr)
    set_drawingml_rpr_element_font(a_rPr, font_name)

def set_pptx_run_font(run, font_name):
    if not run or not font_name:
        return
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        r = run._r
        if hasattr(r, 'get_or_add_rPr'):
            rPr = r.get_or_add_rPr()
        else:
            rPr = r.find(f'{{{A_NS}}}rPr')
            if rPr is None:
                rPr = OxmlElement('a:rPr')
                r.insert(0, rPr)
        set_drawingml_rpr_element_font(rPr, font_name)
    except Exception:
        return

def set_pptx_paragraph_default_font(paragraph, font_name):
    if not paragraph or not font_name:
        return
    try:
        p = paragraph._p
        pPr = p.find(f'{{{A_NS}}}pPr')
        if pPr is None:
            pPr = OxmlElement('a:pPr')
            p.insert(0, pPr)
        defRPr = pPr.find(f'{{{A_NS}}}defRPr')
        if defRPr is None:
            defRPr = OxmlElement('a:defRPr')
            pPr.append(defRPr)
        set_drawingml_rpr_element_font(defRPr, font_name)
    except Exception:
        return

def load_revision_dict(file_path="revision.md", silent=False):
    """
    自动加载校准词典 (revision.md)
    """
    mapping = {}
    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if not line or line.startswith("#"): continue
                    if not re.match(r'^[-*+]\s*', line): continue
                    content = re.sub(r'^[-*+]\s*', '', line)
                    if "格式" in content: continue
                    
                    if ":" in content:
                        parts = content.split(":", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
                    elif "：" in content:
                        parts = content.split("：", 1)
                        if len(parts) == 2:
                            err, corr = parts
                            if err.strip(): mapping[err.strip()] = corr.strip()
            if not silent:
                print(f"[系统] 自动加载校准文件成功，共 {len(mapping)} 条有效规则。")
        except Exception as e:
            print(f"[错误] 加载校准文件失败: {e}")
    return mapping

def apply_revisions(text):
    """
    在翻译结果上应用校准映射
    """
    if not text or not revision_map: return text
    sorted_keys = sorted(revision_map.keys(), key=len, reverse=True)
    result_text = text
    for err in sorted_keys:
        if err in result_text:
            result_text = result_text.replace(err, revision_map[err])
    return result_text

def _lang_name(code):
    return LANGUAGE_NAME_MAP.get(code, code)

def _build_qwen_messages(text, from_lang, to_lang, force_no_hangul=False):
    system_prompt = (
        "You are a professional translation engine.\n"
        "Return ONLY the translation.\n"
        "Do not add explanations, notes, quotes, markdown, or extra text.\n"
        "Preserve original line breaks exactly.\n"
        "Preserve punctuation, numbers, placeholders, and symbols as much as possible."
    )
    user_prompt = (
        f"Translate the following text from {_lang_name(from_lang)} to {_lang_name(to_lang)}.\n"
        + ("Do NOT leave any Korean characters in the output.\n" if force_no_hangul else "")
        + "Text:\n"
        + str(text)
    )
    return [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]

def _ollama_translate_one(text, from_lang, to_lang, force_no_hangul=False):
    base_url, model_name = _check_ollama_runtime()
    messages = _build_qwen_messages(text, from_lang, to_lang, force_no_hangul=force_no_hangul)
    payload = {
        "model": model_name,
        "messages": messages,
        "stream": False,
        "options": {
            "temperature": 0,
            "num_predict": int(QWEN_DEFAULT_MAX_NEW_TOKENS),
        },
    }
    # #region debug-point E:ollama-chat-request
    _debug_report("E", "_ollama_translate_one:request", "发送 Ollama 请求", {"base_url": base_url, "model": model_name, "from_lang": from_lang, "to_lang": to_lang})
    # #endregion
    resp = requests.post(
        f"{base_url}/api/chat",
        json=payload,
        timeout=max(30, int(OLLAMA_TIMEOUT)),
    )
    resp.raise_for_status()
    body = resp.json()
    message = body.get("message") or {}
    text_out = str(message.get("content", "") or "").strip()
    # #region debug-point E:ollama-chat-response
    _debug_report("E", "_ollama_translate_one:response", "收到 Ollama 响应", {"model": model_name, "output_len": len(text_out)})
    # #endregion
    return text_out


def _qwen_generate_batch(text_list, from_lang, to_lang, force_no_hangul=False):
    if not text_list:
        return []

    # #region debug-point E:generate-entry
    _debug_report("E", "_qwen_generate_batch:entry", "进入批量生成", {"batch_size": len(text_list), "from_lang": from_lang, "to_lang": to_lang, "force_no_hangul": force_no_hangul})
    # #endregion
    worker_count = max(1, min(int(OLLAMA_MAX_WORKERS), len(text_list), 3))
    outputs = [None] * len(text_list)

    # 小并发请求 Ollama，CPU 机器默认限制在 2-3 路更稳。
    if worker_count == 1:
        for idx, text in enumerate(text_list):
            outputs[idx] = _ollama_translate_one(text, from_lang, to_lang, force_no_hangul=force_no_hangul)
    else:
        with ThreadPoolExecutor(max_workers=worker_count) as executor:
            future_to_idx = {
                executor.submit(_ollama_translate_one, text, from_lang, to_lang, force_no_hangul): idx
                for idx, text in enumerate(text_list)
            }
            for future in as_completed(future_to_idx):
                idx = future_to_idx[future]
                outputs[idx] = future.result()

    _mark_api(f"ollama_batch={len(text_list)}")
    return outputs

def qwen_translate(q, from_lang, to_lang):
    q = _normalize_text_newlines(q)
    if not q or not q.strip():
        return q
    cache_key = (from_lang, to_lang, str(q))
    cached = _translation_cache.get(cache_key)
    if cached is not None:
        return cached

    original_stripped = str(q).strip()
    strict_no_hangul = (from_lang == "kor") and (to_lang != "kor") and bool(HANGUL_RE.search(original_stripped))
    translated = _qwen_generate_batch([q], from_lang, to_lang, force_no_hangul=False)[0]
    if strict_no_hangul and (translated.strip() == original_stripped or HANGUL_RE.search(translated)):
        translated = _qwen_generate_batch([q], from_lang, to_lang, force_no_hangul=True)[0]
    if not translated.strip():
        translated = q
    _translation_cache[cache_key] = translated
    return translated

def _normalize_text_newlines(text):
    if text is None:
        return ""
    return str(text).replace("\r\n", "\n").replace("\r", "\n")

def _env_int(name, default):
    try:
        v = os.getenv(name)
        if v is None or str(v).strip() == "":
            return int(default)
        return int(str(v).strip())
    except Exception:
        return int(default)

def _env_float(name, default):
    try:
        v = os.getenv(name)
        if v is None or str(v).strip() == "":
            return float(default)
        return float(str(v).strip())
    except Exception:
        return float(default)

def _fallback_translate_one(text, from_lang, to_lang):
    text = _normalize_text_newlines(text)
    if not text or not text.strip():
        return text
    if "\n" not in text:
        return qwen_translate(text, from_lang, to_lang)
    lines = text.split("\n")
    out_lines = []
    for line in lines:
        if line.strip():
            out_lines.append(qwen_translate(line, from_lang, to_lang))
        else:
            out_lines.append(line)
    return "\n".join(out_lines)

def qwen_translate_batch(text_list, from_lang, to_lang):
    if not text_list:
        return []
    results = [None] * len(text_list)

    pending_indices = []
    pending_texts = []
    for i, text in enumerate(text_list):
        if text is None:
            results[i] = text
            continue
        text_s = str(text)
        if not text_s.strip():
            results[i] = text_s
            continue
        
        # --- [新增] 智能跳过逻辑 ---
        if should_skip_translation(text_s, to_lang):
            results[i] = text_s
            continue
            
        cache_key = (from_lang, to_lang, text_s)
        cached = _translation_cache.get(cache_key)
        if cached is not None:
            results[i] = cached
            continue
        pending_indices.append(i)
        pending_texts.append(text_s)

    if not pending_texts:
        return results

    batch_size = max(1, _env_int("QWEN_BATCH_SIZE", _qwen_batch_size))
    total_chunks = (len(pending_texts) + batch_size - 1) // batch_size
    print(f"[系统] 准备处理 {total_chunks} 个 Ollama 翻译批次 (批量大小: {batch_size})...")

    for chunk_idx, start in enumerate(range(0, len(pending_texts), batch_size)):
        chunk_texts = pending_texts[start:start + batch_size]
        chunk_indices = pending_indices[start:start + batch_size]
        msg = f"正在请求批次 {chunk_idx + 1}/{total_chunks} ({len(chunk_texts)} 条文本)..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        try:
            translated_chunk = _qwen_generate_batch(chunk_texts, from_lang, to_lang, force_no_hangul=False)
            for idx_in_orig, translated in zip(chunk_indices, translated_chunk):
                if translated is not None:
                    original_stripped = str(text_list[idx_in_orig] or "").strip()
                    if (from_lang == "kor") and (to_lang != "kor") and (translated.strip() == original_stripped or HANGUL_RE.search(translated)):
                        translated = _fallback_translate_one(text_list[idx_in_orig], from_lang, to_lang)
                    results[idx_in_orig] = translated
                    cache_key = (from_lang, to_lang, str(text_list[idx_in_orig]))
                    _translation_cache[cache_key] = translated
        except Exception as e:
            print(f"[错误] 批次 {chunk_idx + 1} 失败，自动降级逐条处理: {e}")
            for idx_in_orig, original_text in zip(chunk_indices, chunk_texts):
                translated = _fallback_translate_one(original_text, from_lang, to_lang)
                results[idx_in_orig] = translated
                cache_key = (from_lang, to_lang, str(text_list[idx_in_orig]))
                _translation_cache[cache_key] = translated

    for i in range(len(results)):
        if results[i] is None:
            results[i] = _fallback_translate_one(text_list[i], from_lang, to_lang)
    return results

def get_translation_batch(text_list):
    if not text_list:
        return []
    direction = translation_direction.get()
    from_lang = FROM_LANG_MAP.get(direction, 'auto')
    to_lang = TO_LANG_MAP.get(direction, 'auto')

    translated_list = qwen_translate_batch(text_list, from_lang, to_lang)
    out = []
    for t in translated_list:
        out.append(apply_revisions(t))
    return out

def get_translation(text):
    if not text or not text.strip(): return text
    direction = translation_direction.get()
    from_lang = FROM_LANG_MAP.get(direction, 'auto')
    to_lang = TO_LANG_MAP.get(direction, 'auto')
    
    # --- [新增] 智能跳过逻辑 ---
    if should_skip_translation(text, to_lang):
        return text
    
    translated_text = qwen_translate(text, from_lang, to_lang)
    # 事后校准
    translated_text = apply_revisions(translated_text)
    return translated_text

def _ppt_normalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\r\n", "\n").replace("\r", "\n").replace("\v", "\n")

def _ppt_denormalize_linebreaks(text):
    if not text:
        return text
    return str(text).replace("\n", "\v")

def translate_ppt_paragraph(paragraph):
    full_text = paragraph.text
    if not full_text or not str(full_text).strip():
        return

    normalized = _ppt_normalize_linebreaks(full_text)
    translated = get_translation(normalized)
    original_texts.append(normalized)
    translated_texts.append(translated)
    _apply_ppt_paragraph_translation(paragraph, normalized, translated)

def _apply_ppt_paragraph_translation(paragraph, normalized, translated):
    if append_translation.get():
        result_norm = append_translation_to_original(normalized, translated)
    else:
        result_norm = translated

    result_text = _ppt_denormalize_linebreaks(result_norm)

    if paragraph.runs:
        first_run = paragraph.runs[0]
    else:
        first_run = paragraph.add_run()

    original_font_name = first_run.font.name
    original_size = first_run.font.size
    original_bold = first_run.font.bold
    original_italic = first_run.font.italic
    original_underline = first_run.font.underline

    first_run.text = result_text
    for r in paragraph.runs[1:]:
        r.text = ""

    target_font_name = get_target_font_name()
    if target_font_name:
        set_pptx_run_font(first_run, target_font_name)
        set_pptx_paragraph_default_font(paragraph, target_font_name)
    elif original_font_name:
        first_run.font.name = original_font_name
        first_run.font.size = original_size
        first_run.font.bold = original_bold
        first_run.font.italic = original_italic
        first_run.font.underline = original_underline

def append_translation_to_original(text, translated_text, cell=None):
    text = text.strip()
    translated_text = translated_text.strip()
    result = f"{text}\n{translated_text}" if text and translated_text else (text or translated_text)
    if cell:
        cell.alignment = Alignment(wrap_text=True, vertical='center')
        ws = cell.parent
        row_num = cell.row
        original_height = ws.row_dimensions[row_num].height
        ws.row_dimensions[row_num].height = (original_height * 2) if original_height else 30
    return result

# ==========================================
# 文件处理逻辑 (完全基于 V2.9 的稳定代码)
# ==========================================

def translate_shape_for_ppt(shape):
    if shape.has_text_frame:
        paragraphs = list(shape.text_frame.paragraphs)
        _translate_ppt_paragraphs_batch(paragraphs)
    elif shape.shape_type == 6:  # 组合形状
        for sub_shape in shape.shapes:
            translate_shape_for_ppt(sub_shape)
    elif shape.has_table:  # 处理表格形状
        table = shape.table
        paragraphs = []
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraphs.append(paragraph)
        _translate_ppt_paragraphs_batch(paragraphs)

def _translate_ppt_paragraphs_batch(paragraphs):
    if not paragraphs:
        return
    jobs = []
    texts = []
    for paragraph in paragraphs:
        full_text = paragraph.text
        if not full_text or not str(full_text).strip():
            continue
        normalized = _ppt_normalize_linebreaks(full_text)
        jobs.append((paragraph, normalized))
        texts.append(normalized)
    if not jobs:
        return
    translated_list = get_translation_batch(texts)
    for (paragraph, normalized), translated in zip(jobs, translated_list):
        original_texts.append(normalized)
        translated_texts.append(translated)
        _apply_ppt_paragraph_translation(paragraph, normalized, translated)

def update_ui_status(msg):
    """线程安全地更新 UI 状态"""
    _mark_progress(msg)
    try:
        if root and root.winfo_exists():
            root.after(0, lambda: status_label.config(text=msg))
    except Exception:
        pass

def translate_ppt(input_file, output_file):
    prs = Presentation(input_file)
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        msg = f"正在翻译 PPT: 第 {i}/{total_slides} 页..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        for shape in slide.shapes:
            translate_shape_for_ppt(shape)
    
    print(f"[系统] 正在保存 PPT 文件...")
    prs.save(output_file)
    save_to_corpus(original_texts, translated_texts)
    print(f"[完成] 保存修改后的 PPT 文件完成!")

def clean_sheet_name(name):
    if not name: return "Sheet"
    invalid_chars = r'[\\/?:*\[\](){}<>|"\']'
    clean_name = re.sub(invalid_chars, '', name)
    return clean_name[:31]

def translate_excel_xlsx(input_file, output_file):
    """处理.xlsx格式的Excel文件"""
    wb = load_workbook(input_file, keep_vba=True)
    existing_sheet_names = set()
    original_sheet_names = list(wb.sheetnames)
    translated_sheet_names = get_translation_batch(original_sheet_names) if original_sheet_names else []
    total_sheets = len(original_sheet_names)
    
    for i, sheet_name in enumerate(original_sheet_names, 1):
        msg = f"正在翻译 Excel: 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        translated_sheet_name = translated_sheet_names[i - 1] if (i - 1) < len(translated_sheet_names) else get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        ws = wb[sheet_name]
        ws.title = unique_sheet_name
        
        # 准备并发任务池
        all_row_jobs = []
        for row_idx, row in enumerate(ws.iter_rows(), 1):
            row_data = []
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    original_value = str(cell.value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    row_data.append({
                        'cell': cell,
                        'orig': normalized_value,
                        'font': Font(**cell.font.__dict__),
                        'border': Border(**cell.border.__dict__),
                        'align': Alignment(**cell.alignment.__dict__),
                        'fill': PatternFill(**cell.fill.__dict__)
                    })
            if row_data:
                all_row_jobs.append((row_idx, row_data))

        if not all_row_jobs:
            continue

        print(f"[系统] 工作表 {unique_sheet_name} 开始并发翻译 ({len(all_row_jobs)} 行有内容)")
        
        def process_single_row(row_info):
            # 检查 UI 是否还在
            try:
                if not root or not root.winfo_exists():
                    return row_info[0], row_info[1], [j['orig'] for j in row_info[1]]
            except Exception:
                return row_info[0], row_info[1], [j['orig'] for j in row_info[1]]

            r_idx, r_jobs = row_info
            texts_to_translate = [j['orig'] for j in r_jobs]
            try:
                translated_texts_list = get_translation_batch(texts_to_translate)
                return r_idx, r_jobs, translated_texts_list
            except Exception as e:
                print(f"[错误] 行 {r_idx} 翻译失败: {e}")
                return r_idx, r_jobs, [j['orig'] for j in r_jobs]

        with ThreadPoolExecutor(max_workers=_kimi_max_workers) as executor:
            future_to_row = {executor.submit(process_single_row, job): job for job in all_row_jobs}
            completed_count = 0
            row_results = {}  # 用于按顺序存储结果
            for future in as_completed(future_to_row):
                r_idx, r_jobs, translated_list = future.result()
                row_results[r_idx] = (r_jobs, translated_list)
                completed_count += 1
                
                if completed_count % 5 == 0 or completed_count == len(all_row_jobs):
                    row_msg = f"并发处理中: {unique_sheet_name} 已完成 {completed_count}/{len(all_row_jobs)} 行..."
                    update_ui_status(row_msg)
                    if completed_count % 20 == 0:
                        print(f"[进度] {row_msg}")

            # 所有线程完成后，按 r_idx 顺序写入单元格和语料库
            for r_idx in sorted(row_results.keys()):
                r_jobs, translated_list = row_results[r_idx]
                for job, trans_txt in zip(r_jobs, translated_list):
                    original_texts.append(job['orig'])
                    translated_texts.append(trans_txt)
                    
                    if append_translation.get():
                        job['cell'].value = append_translation_to_original(job['orig'], trans_txt, job['cell'])
                    else:
                        job['cell'].value = trans_txt
                    
                    target_font_name = get_target_font_name()
                    if target_font_name:
                        job['cell'].font = Font(
                            name=target_font_name,
                            size=job['font'].size,
                            bold=job['font'].bold,
                            italic=job['font'].italic,
                            underline=job['font'].underline,
                            color=job['font'].color
                        )
                    else:
                        job['cell'].font = job['font']
                    
                    job['cell'].border = job['border']
                    job['cell'].alignment = job['align']
                    job['cell'].fill = job['fill']
                    job['cell'].alignment = Alignment(wrap_text=True, vertical='center')

    wb.save(output_file)
    save_to_corpus(original_texts, translated_texts)

def translate_excel_xls(input_file, output_file):
    """处理.xls格式的Excel文件，使用pandas库，保存为.xlsx格式"""
    base_output = os.path.splitext(output_file)[0]
    output_file_xlsx = base_output + '.xlsx'
    target_font_name = get_target_font_name()
    excel_file = pd.ExcelFile(input_file)
    writer = pd.ExcelWriter(output_file_xlsx, engine='openpyxl')
    existing_sheet_names = set()
    original_sheet_names = list(excel_file.sheet_names)
    translated_sheet_names = get_translation_batch(original_sheet_names) if original_sheet_names else []
    total_sheets = len(original_sheet_names)
    
    for i, sheet_name in enumerate(original_sheet_names, 1):
        msg = f"正在翻译 Excel(.xls): 第 {i}/{total_sheets} 个工作表 ({sheet_name})..."
        print(f"[进度] {msg}")
        update_ui_status(msg)
        
        translated_sheet_name = translated_sheet_names[i - 1] if (i - 1) < len(translated_sheet_names) else get_translation(sheet_name)
        cleaned_sheet_name = clean_sheet_name(translated_sheet_name)
        
        unique_sheet_name = cleaned_sheet_name
        count = 1
        while unique_sheet_name in existing_sheet_names:
            temp_name = f"{cleaned_sheet_name[:27]}_{count}"
            unique_sheet_name = temp_name[:31]
            count += 1
        existing_sheet_names.add(unique_sheet_name)
        
        df = pd.read_excel(input_file, sheet_name=sheet_name)
        total_rows = len(df.index)
        
        # 并发处理行
        all_row_jobs = []
        for idx_num, idx in enumerate(df.index, 1):
            row_texts = []
            row_cols = []
            for col in df.columns:
                cell_value = df.at[idx, col]
                if pd.notna(cell_value) and isinstance(cell_value, str):
                    original_value = str(cell_value)
                    normalized_value = original_value.replace("\r\n", "\n").replace("\r", "\n")
                    row_texts.append(normalized_value)
                    row_cols.append(col)
            if row_texts:
                all_row_jobs.append((idx, row_texts, row_cols))

        if all_row_jobs:
            print(f"[系统] 工作表 {unique_sheet_name} 开始并发翻译 ({len(all_row_jobs)} 行有内容)")
            
            def process_xls_row(row_info):
                # 检查 UI 是否还在
                try:
                    if not root or not root.winfo_exists():
                        return row_info[0], row_info[1], row_info[2], row_info[1]
                except Exception:
                    return row_info[0], row_info[1], row_info[2], row_info[1]

                idx, r_texts, r_cols = row_info
                try:
                    translated_list = get_translation_batch(r_texts)
                    return idx, r_texts, r_cols, translated_list
                except Exception as e:
                    print(f"[错误] 行 {idx} 翻译失败: {e}")
                    return idx, r_texts, r_cols, r_texts

            with ThreadPoolExecutor(max_workers=_kimi_max_workers) as executor:
                future_to_row = {executor.submit(process_xls_row, job): job for job in all_row_jobs}
                completed_count = 0
                row_results = {}
                for future in as_completed(future_to_row):
                    idx, r_texts, r_cols, translated_list = future.result()
                    row_results[idx] = (r_texts, r_cols, translated_list)
                    completed_count += 1
                    
                    if completed_count % 10 == 0 or completed_count == len(all_row_jobs):
                        row_msg = f"并发处理中: {unique_sheet_name} 已完成 {completed_count}/{len(all_row_jobs)} 行..."
                        update_ui_status(row_msg)

                # 按顺序写回和记录
                for idx in sorted(row_results.keys()):
                    r_texts, r_cols, translated_list = row_results[idx]
                    for orig, col, trans in zip(r_texts, r_cols, translated_list):
                        original_texts.append(orig)
                        translated_texts.append(trans)
                        df.at[idx, col] = append_translation_to_original(orig, trans) if append_translation.get() else trans

        df.to_excel(writer, sheet_name=unique_sheet_name, index=False)
    
    writer.close()
    if target_font_name:
        try:
            wb = load_workbook(output_file_xlsx)
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            cell.font = Font(
                                name=target_font_name,
                                size=cell.font.size,
                                bold=cell.font.bold,
                                italic=cell.font.italic,
                                underline=cell.font.underline,
                                color=cell.font.color
                            )
            wb.save(output_file_xlsx)
        except Exception: pass
    save_to_corpus(original_texts, translated_texts)
    return output_file_xlsx

def translate_word(input_file, output_file):
    doc = Document(input_file)
    target_font_name = get_target_font_name()
    processed_elements = set()

    def set_style_font(style_element, font_name):
        if style_element is None or not font_name:
            return
        rPr = style_element.find(qn('w:rPr'))
        if rPr is None:
            rPr = OxmlElement('w:rPr')
            style_element.append(rPr)
        rFonts = rPr.find(qn('w:rFonts'))
        if rFonts is None:
            rFonts = OxmlElement('w:rFonts')
            rPr.append(rFonts)
        rFonts.set(qn('w:ascii'), font_name)
        rFonts.set(qn('w:hAnsi'), font_name)
        rFonts.set(qn('w:eastAsia'), font_name)
        rFonts.set(qn('w:cs'), font_name)
        for k in (qn('w:asciiTheme'), qn('w:hAnsiTheme'), qn('w:eastAsiaTheme'), qn('w:csTheme')):
            try:
                rFonts.attrib.pop(k, None)
            except Exception:
                pass

    # --- 阶段 1: 扫描并收集所有需要翻译的文本 (用于高性能批量请求) ---
    print(f"[系统] 正在预扫描 Word 提取待翻译文本...")
    pending_texts = []
    
    def collect_texts_from_container(container_el):
        # 1. 扫描所有标准段落 <w:p>
        for p_el in container_el.xpath('.//w:p'):
            try:
                full_text = Paragraph(p_el, doc).text
                if full_text and full_text.strip():
                    # 兼容协议中的软回车 \v
                    for line in full_text.replace('\v', '\n').split('\n'):
                        if line.strip(): pending_texts.append(line.strip())
            except Exception: continue
        # 2. 扫描图形文字 <a:p>
        for p_el in container_el.xpath('.//a:p'):
            try:
                t_nodes = p_el.xpath('.//a:t')
                combined = "".join([n.text for n in t_nodes if n.text])
                if combined.strip(): pending_texts.append(combined.strip())
            except Exception: continue

    collect_texts_from_container(doc.element.body)
    for section in doc.sections:
        for h in [section.header, section.first_page_header, section.even_page_header]:
            if h: collect_texts_from_container(h._element)
        for f in [section.footer, section.first_page_footer, section.even_page_footer]:
            if f: collect_texts_from_container(f._element)

    # --- 阶段 2: 统一高性能批量翻译 ---
    if pending_texts:
        unique_texts = list(dict.fromkeys(pending_texts)) # 去重并保持顺序
        try:
            from_lang, to_lang = translation_direction.get().split('2')
        except Exception:
            from_lang, to_lang = "ko", "zh"
            
        # 仅翻译不在缓存中的文本
        to_translate = [t for t in unique_texts if _translation_cache.get((from_lang, to_lang, t)) is None]
        
        if to_translate:
            print(f"[系统] 发现 {len(to_translate)} 条新文本，正在发起并发批量请求...")
            get_translation_batch(to_translate) # 该函数会自动填充 _translation_cache
        else:
            print(f"[系统] 所有文本已在缓存中，直接进行回填。")

    # --- 阶段 3: 第二次扫描并回填翻译 (由于有缓存，此步瞬间完成) ---
    def translate_word_paragraph(paragraph):
        """核心段落翻译：支持软回车拆分、行对行对照、去重处理"""
        p_id = id(paragraph._element)
        if p_id in processed_elements:
            return
        processed_elements.add(p_id)

        full_text = paragraph.text
        if not full_text or not str(full_text).strip():
            return

        # --- 捕捉原句的模板 Run (用于对照模式下的格式克隆) ---
        template_run = paragraph.runs[0] if paragraph.runs else None

        # 使用 splitlines() 稳健处理软回车
        lines = [l.strip() for l in full_text.replace('\v', '\n').splitlines() if l.strip()]
        if not lines:
            return

        print(f"[回填] Word 段落 ({len(lines)} 行): {lines[0][:30]}...")

        if append_translation.get():
            # 【对照模式：精准行对行回填】
            translated_lines = [get_translation(l) for l in lines] # 此时 get_translation 走缓存
            
            for orig, trans in zip(lines, translated_lines):
                original_texts.append(orig)
                translated_texts.append(trans)
            
            # 清空段落所有内容
            p_el = paragraph._element
            for r in p_el.xpath('./w:r'):
                p_el.remove(r)

            # 按顺序填充原文和译文
            for i, (orig, trans) in enumerate(zip(lines, translated_lines)):
                # 1. 写入原文 (应用模板格式)
                run_orig = paragraph.add_run(orig)
                if template_run:
                    apply_run_format(template_run, run_orig)
                
                run_orig.add_break(WD_BREAK.LINE)
                
                # 2. 写入译文 (应用模板格式，并覆盖目标字体)
                run_trans = paragraph.add_run(trans)
                if template_run:
                    apply_run_format(template_run, run_trans)
                
                if target_font_name:
                    set_docx_run_font(run_trans, target_font_name)
                
                # 组间换行
                if i < len(lines) - 1: run_trans.add_break(WD_BREAK.LINE)
        else:
            # 【替换模式】
            normalized = "\n".join(lines)
            t = get_translation(normalized)
            original_texts.append(normalized)
            translated_texts.append(t)
            
            if paragraph.runs:
                first_run = paragraph.runs[0]
                first_run.text = t
                for r in paragraph.runs[1:]:
                    if not bool(xpath_with_ns(r.element, './/w:drawing')): r.text = ""
                if target_font_name:
                    set_docx_run_font(first_run, target_font_name)
            else:
                new_run = paragraph.add_run(t)
                if target_font_name:
                    set_docx_run_font(new_run, target_font_name)

    def fill_translations_in_container(container_el):
        # 1. 回填段落
        for p_el in container_el.xpath('.//w:p'):
            try: translate_word_paragraph(Paragraph(p_el, doc))
            except Exception: continue
        # 2. 回填图形
        for p_el in container_el.xpath('.//a:p'):
            p_id = id(p_el)
            if p_id in processed_elements: continue
            processed_elements.add(p_id)
            t_nodes = p_el.xpath('.//a:t')
            if not t_nodes: continue
            combined = "".join([n.text for n in t_nodes if n.text])
            if not combined.strip(): continue
            t = get_translation(combined) # 走缓存
            original_texts.append(combined)
            translated_texts.append(t)
            if append_translation.get():
                t_nodes[0].text = combined
                for n in t_nodes[1:]: n.text = ""
                r_nodes = p_el.xpath('.//a:r')
                if r_nodes:
                    last_r = r_nodes[0]
                    br = OxmlElement('a:br')
                    last_r.addnext(br)
                    new_r = OxmlElement('a:r')
                    new_t = OxmlElement('a:t')
                    new_t.text = t
                    new_r.append(new_t)
                    br.addnext(new_r)
                    if target_font_name: set_drawingml_r_element_font(new_r, target_font_name)
            else:
                t_nodes[0].text = t
                for n in t_nodes[1:]: n.text = ""

    print(f"[系统] 正在回填 Word 翻译内容...")
    fill_translations_in_container(doc.element.body)
    for section in doc.sections:
        for h in [section.header, section.first_page_header, section.even_page_header]:
            if h: fill_translations_in_container(h._element)
        for f in [section.footer, section.first_page_footer, section.even_page_footer]:
            if f: fill_translations_in_container(f._element)

    # 字体加固 (仅针对非对照模式，或者作为样式兜底)
    if target_font_name:
        try:
            for s in doc.styles:
                try:
                    if getattr(s, "font", None) is not None: s.font.name = target_font_name
                except Exception: pass
                try: set_style_font(getattr(s, "_element", None), target_font_name)
                except Exception: pass
            
            # 注意：移除对照模式下的全局 w:r 强制刷，防止覆盖原文原字体
            if not append_translation.get():
                for r in doc.element.xpath('.//w:r'): set_docx_r_element_font(r, target_font_name)
            
            for r in doc.element.xpath('.//a:r'): set_drawingml_r_element_font(r, target_font_name)
        except Exception: pass

    doc.save(output_file)
    save_to_corpus(original_texts, translated_texts)

# --- 线程控制 ---

def check_direction_mismatch(input_file):
    """
    预检：如果翻译方向和文件内容明显不匹配，弹出预警
    """
    try:
        direction = translation_direction.get()
        from_lang = FROM_LANG_MAP.get(direction)
        to_lang = TO_LANG_MAP.get(direction)
        
        sample_text = ""
        ext = os.path.splitext(input_file)[1].lower()
        if ext == '.docx':
            doc = Document(input_file)
            # 取前5个段落
            sample_text = "\n".join([p.text for p in doc.paragraphs[:5] if p.text.strip()])
        elif ext == '.xlsx':
            wb = load_workbook(input_file, read_only=True, data_only=True)
            ws = wb.active
            cells = []
            for row in ws.iter_rows(max_row=10):
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        cells.append(str(cell.value))
            sample_text = "\n".join(cells)
        elif ext == '.ppt' or ext == '.pptx':
            prs = Presentation(input_file)
            texts = []
            for slide in prs.slides[:3]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            sample_text = "\n".join(texts)

        if not sample_text.strip():
            return True

        # --- [优化] 更加通用的语言特征预检 ---
        has_korean = bool(HANGUL_RE.search(sample_text))
        has_chinese = bool(CHINESE_RE.search(sample_text))
        has_japanese = bool(re.search(r'[\u3040-\u30ff]', sample_text))
        # 检测连续的英文字母，排除掉零散的符号
        has_english = bool(re.search(r'[a-zA-Z]{5,}', sample_text)) 

        # 构造易读的方向标签
        dir_labels = {
            'ko2zh': '韩 -> 中', 'zh2ko': '中 -> 韩', 'ko2vi': '韩 -> 越', 'ko2en': '韩 -> 英',
            'zh2en': '中 -> 英', 'en2zh': '英 -> 中', 'zh_tw2en': '繁中 -> 英', 'en2zh_tw': '英 -> 繁中',
            'zh2ja': '中 -> 日', 'ja2zh': '日 -> 中', 'en2ko': '英 -> 韩', 'vi2zh': '越 -> 中'
        }
        current_dir_label = dir_labels.get(direction, direction)

        warning_reason = ""
        # 核心预警逻辑：如果检测到的语言 既不是源语言 也不是 目标语言，则报警
        if has_chinese and from_lang != 'zh' and to_lang not in ('zh', 'zh_tw'):
            warning_reason = "中文"
        elif has_korean and from_lang != 'kor' and to_lang != 'kor':
            warning_reason = "韩文"
        elif has_japanese and from_lang != 'ja' and to_lang != 'ja':
            warning_reason = "日文"
        elif has_english and from_lang != 'en' and to_lang != 'en':
            # 英文报警稍微严谨一点：只有当本该出现的源语言完全没出现时才报警
            if from_lang == 'kor' and not has_korean: warning_reason = "英文"
            elif from_lang == 'zh' and not has_chinese: warning_reason = "英文"
            elif from_lang == 'ja' and not has_japanese: warning_reason = "英文"

        if warning_reason:
            return messagebox.askyesno("预警", f"检测到当前翻译方向为 [{current_dir_label}]，但文件内容似乎是 {warning_reason}（既非源语言也非目标语言）。\n\n是否继续执行翻译？")
        
        return True
    except Exception as e:
        print(f"[预检] 预检过程出错: {e}")
        return True

def start_translation():
    input_file = input_file_entry.get()
    output_folder = output_folder_entry.get()
    ollama_base_url = ollama_base_url_entry.get().strip()
    ollama_model_name = ollama_model_entry.get().strip()
    if not (input_file and output_folder and ollama_base_url and ollama_model_name):
        messagebox.showwarning("提示", "请完整选择输入文件、输出目录、Ollama 服务地址和模型名")
        return
    
    # --- [新增] 翻译方向预检 ---
    if not check_direction_mismatch(input_file):
        return
    
    translate_button.config(state=tk.DISABLED, text="🚀 正在翻译，请稍候...")
    status_label.config(text="任务已启动，请查看终端进度...", foreground="#2980b9")
    thread = threading.Thread(target=run_translation_task, args=(input_file, output_folder))
    thread.daemon = True
    thread.start()

def run_translation_task(input_file, output_folder):
    _watchdog_start()
    start_time = time.time()
    try:
        # #region debug-point B:translation-thread-entry
        _debug_report("B", "run_translation_task:entry", "翻译线程启动", {"input_file": input_file, "output_folder": output_folder, "thread_name": threading.current_thread().name})
        # #endregion
        global revision_map, OLLAMA_BASE_URL, OLLAMA_MODEL_NAME
        file_ext = os.path.splitext(input_file)[1].lower()
        custom_name = custom_filename_entry.get().strip()
        ollama_base_url = ollama_base_url_entry.get().strip()
        ollama_model_name = ollama_model_entry.get().strip()
        if ollama_base_url:
            OLLAMA_BASE_URL = _normalize_ollama_base_url(ollama_base_url)
        if ollama_model_name:
            OLLAMA_MODEL_NAME = ollama_model_name
        output_file = os.path.join(output_folder, (custom_name if custom_name else f"translated_v4.0_{os.path.basename(input_file).split('.')[0]}") + file_ext)
        
        # --- [V2.14 增强功能：物理克隆以保留图片和绘图] ---
        # 即使 V2.9 也不支持 Excel 绘图保留，这里通过物理复制尝试最大化兼容性
        try:
            if os.path.exists(output_file):
                os.remove(output_file) # 尝试删除旧文件，如果被占用会在这里报错
            shutil.copy2(input_file, output_file)
        except PermissionError:
            raise Exception(f"目标文件已被占用，请先关闭 Excel/Word/PPT: {os.path.basename(output_file)}")
        
        original_texts.clear()
        translated_texts.clear()
        revision_map = load_revision_dict("revision.md", silent=True)
        print(f"\n[系统] 开始翻译任务: {os.path.basename(input_file)}")
        print(f"[系统] Ollama 服务地址: {OLLAMA_BASE_URL}")
        print(f"[系统] Ollama 模型名: {OLLAMA_MODEL_NAME}")
        print(f"[系统] 校准规则加载: {len(revision_map)} 条")
        update_ui_status(f"开始翻译: {os.path.basename(input_file)}")
        
        # 统一使用 output_file 作为操作对象，实现“原地翻译”
        if file_ext in ['.ppt', '.pptx']: translate_ppt(output_file, output_file)
        elif file_ext == '.xlsx': translate_excel_xlsx(output_file, output_file)
        elif file_ext == '.xls': 
            # .xls 比较特殊，必须另存为 .xlsx
            output_file = translate_excel_xls(input_file, output_file)
        elif file_ext == '.docx': translate_word(output_file, output_file)
        
        end_time = time.time()
        duration_minutes = (end_time - start_time) / 60
        print(f"[完成] 文件已保存至: {output_file}")
        print(f"[统计] 翻译总耗时: {duration_minutes:.2f} 分钟\n")
        root.after(0, lambda: translation_done_callback(output_file, duration_minutes))
    except Exception as e:
        err_msg = str(e)
        # #region debug-point B:translation-thread-exception
        _debug_report("B", "run_translation_task:exception", "翻译线程捕获异常", {"error": err_msg, "traceback": traceback.format_exc()[:4000]})
        # #endregion
        print(f"[错误] 详情: {err_msg}")
        root.after(0, lambda: translation_failed_callback(err_msg))
    finally:
        # #region debug-point B:translation-thread-finally
        _debug_report("B", "run_translation_task:finally", "翻译线程进入 finally", {"input_file": input_file})
        # #endregion
        _watchdog_stop()

def translation_done_callback(output_file, duration_minutes):
    translate_button.config(state=tk.NORMAL, text="🚀 开始 Qwen2.5 (Ollama/GGUF) 翻译")
    status_label.config(text=f"翻译任务已圆满完成！(用时 {duration_minutes:.2f} 分钟)", foreground="#27ae60")
    messagebox.showinfo("成功", f"翻译完成！\n用时：{duration_minutes:.2f} 分钟\n文件保存至：{output_file}")

def translation_failed_callback(error_msg):
    translate_button.config(state=tk.NORMAL, text="🚀 开始 Qwen2.5 (Ollama/GGUF) 翻译")
    status_label.config(text=f"翻译过程出错", foreground="#e74c3c")
    messagebox.showerror("错误", f"发生异常：{error_msg}")

def save_to_corpus(orig, trans):
    if generate_corpus.get() and orig:
        if not os.path.exists('Corpus'): os.makedirs('Corpus')
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        try:
            direction = translation_direction.get()
        except Exception:
            direction = "unknown"
        corpus_file = f'Corpus/Corpus_v4.0_{direction}_{timestamp}.xlsx'
        to_lang = TO_LANG_MAP.get(direction)
        seen = set()
        filtered_orig = []
        filtered_trans = []
        for o, t in zip(orig, trans):
            if o is None or t is None:
                continue
            o = str(o).replace("\r\n", "\n").replace("\r", "\n").strip()
            t = str(t).replace("\r\n", "\n").replace("\r", "\n").strip()
            if not o or not t:
                continue

            # 3. 原文=译文 不记录
            if o == t:
                continue

            # 检查语言特征
            has_korean = bool(HANGUL_RE.search(o))
            has_chinese = bool(CHINESE_RE.search(o))

            # 1. 核心准则：只要含有韩文就要记录 (满足 "原文只要含有韩文就要记录")
            if has_korean:
                pass 
            else:
                # 如果不含韩文：
                # A. 如果包含中文，且目标是中文 (针对韩翻中方向) -> 不记录 (即 "纯中文不录")
                # 注意：如果包含中文但目标不是中文 (如中翻韩、中翻英)，则记录
                if has_chinese and to_lang in ("zh", "zh_tw"):
                    continue
                
                # B. 如果既不含韩文也不含中文 -> 说明只有数字、英文、符号的组合 -> 不记录
                # 这完全符合 "数字/英文/符号 只有全部是的时候才能不记录"
                if not has_chinese:
                    continue

            # 经过筛选后，执行去重检查
            key = o
            if key in seen:
                continue
            seen.add(key)
            filtered_orig.append(o)
            filtered_trans.append(t)

        if not filtered_orig:
            return

        pd.DataFrame(
            {'序号': range(1, len(filtered_orig) + 1), '翻译前': filtered_orig, '翻译后': filtered_trans}
        ).to_excel(corpus_file, index=False)

# --- UI 布局 ---
root = tk.Tk()
root.title("Qwen2.5 (Ollama/GGUF) 翻译工具 V4.0")
root.geometry("700x850")
root.configure(bg="#f5f6fa")
style = ttk.Style()
style.theme_use('clam')
style.configure("TFrame", background="#f5f6fa")
style.configure("TLabel", background="#f5f6fa", font=("微软雅黑", 10))
style.configure("Header.TLabel", font=("微软雅黑", 14, "bold"), foreground="#166534")
main_frame = ttk.Frame(root, padding="20")
main_frame.pack(fill="both", expand=True)
ttk.Label(main_frame, text="Qwen2.5 (Ollama/GGUF) 翻译 & 自动校准系统", style="Header.TLabel").pack(pady=(0, 20))
file_card = ttk.LabelFrame(main_frame, text=" 文件设置 ", padding=15)
file_card.pack(fill="x", pady=10)
ttk.Label(file_card, text="待翻译文件:").grid(row=0, column=0, sticky="w", pady=5)
input_file_entry = ttk.Entry(file_card, width=50)
input_file_entry.grid(row=0, column=1, padx=10)
ttk.Button(file_card, text="选择文件", command=lambda: (input_file_entry.delete(0, tk.END), input_file_entry.insert(0, filedialog.askopenfilename()))).grid(row=0, column=2)
ttk.Label(file_card, text="保存位置:").grid(row=1, column=0, sticky="w", pady=5)
output_folder_entry = ttk.Entry(file_card, width=50)
output_folder_entry.grid(row=1, column=1, padx=10)
ttk.Button(file_card, text="选择目录", command=lambda: (output_folder_entry.delete(0, tk.END), output_folder_entry.insert(0, filedialog.askdirectory()))).grid(row=1, column=2)
ttk.Label(file_card, text="Ollama 服务地址:").grid(row=2, column=0, sticky="w", pady=5)
ollama_base_url_entry = ttk.Entry(file_card, width=50)
ollama_base_url_entry.grid(row=2, column=1, padx=10)
ollama_base_url_entry.insert(0, OLLAMA_BASE_URL)
ttk.Label(file_card, text="Ollama 模型名:").grid(row=3, column=0, sticky="w", pady=5)
ollama_model_entry = ttk.Entry(file_card, width=50)
ollama_model_entry.grid(row=3, column=1, padx=10, columnspan=2, sticky="w")
ollama_model_entry.insert(0, OLLAMA_MODEL_NAME)
ttk.Label(file_card, text="自定义文件名:").grid(row=4, column=0, sticky="w", pady=5)
custom_filename_entry = ttk.Entry(file_card, width=50)
custom_filename_entry.grid(row=4, column=1, padx=10, columnspan=2, sticky="w")
dir_card = ttk.LabelFrame(main_frame, text=" 翻译语种 ", padding=15)
dir_card.pack(fill="x", pady=10)
translation_direction = tk.StringVar(value='ko2zh')
lang_grid = ttk.Frame(dir_card)
lang_grid.pack(fill="x")
langs = [("韩 -> 中", 'ko2zh'), ("中 -> 韩", 'zh2ko'), ("韩 -> 越", 'ko2vi'), ("韩 -> 英", 'ko2en'),
         ("中 -> 英", 'zh2en'), ("英 -> 中", 'en2zh'), ("繁中 -> 英", 'zh_tw2en'), ("英 -> 繁中", 'en2zh_tw'),
         ("中 -> 日", 'zh2ja'), ("日 -> 中", 'ja2zh'), ("英 -> 韩", 'en2ko'), ("越 -> 中", 'vi2zh')]
for i, (text, val) in enumerate(langs):
    row, col = i // 4, i % 4
    ttk.Radiobutton(lang_grid, text=text, variable=translation_direction, value=val).grid(row=row, column=col, padx=15, pady=5, sticky="w")
opt_frame = tk.Frame(main_frame, bg="#f5f6fa")
opt_frame.pack(fill="x", pady=10)
append_translation = tk.BooleanVar(value=False)
tk.Checkbutton(opt_frame, text="在原文下方保留翻译对照", variable=append_translation, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
generate_corpus = tk.BooleanVar(value=False)
tk.Checkbutton(opt_frame, text="同步生成语料库 (Corpus)", variable=generate_corpus, bg="#f5f6fa", selectcolor="white").pack(side="left", padx=20)
translate_button = tk.Button(main_frame, text="🚀 开始 Qwen2.5 (Ollama/GGUF) 翻译", command=start_translation, bg="#16a34a", fg="white", font=("微软雅黑", 12, "bold"), relief="flat", cursor="hand2", pady=12)
translate_button.pack(fill="x", pady=20)
status_label = ttk.Label(main_frame, text="就绪：已加载校准词典 (revision.md)", foreground="#7f8c8d")
status_label.pack()

# 初始化加载一次即可
revision_map = load_revision_dict("revision.md")

root.mainloop()
